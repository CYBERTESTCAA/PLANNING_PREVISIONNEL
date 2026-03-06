#!/usr/bin/env python3
"""
Generate two PowerPoint presentations for the Planning Prévisionnel project:
  1) Technical architecture & developer guide
  2) User guide (how to use the application)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color palette ──────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1E, 0x3A, 0x5F)
BLUE_MED   = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1F, 0x1F, 0x1F)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
GREEN      = RGBColor(0x05, 0x96, 0x69)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
RED        = RGBColor(0xDC, 0x26, 0x26)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    """Add a colored rectangle (no border)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_title_bar(slide, color=BLUE_DARK):
    """Add a dark top bar."""
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), color)

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with single-run text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_bullet_slide(slide, items, left=Inches(0.8), top=Inches(2.0), width=Inches(11.5), font_size=15, color=BLACK, spacing=Pt(8)):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        # Handle bold prefix (text before first colon)
        if " : " in item:
            parts = item.split(" : ", 1)
            run = p.add_run()
            run.text = "•  " + parts[0] + " : "
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = color
            run.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = "•  " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return txBox

def add_card(slide, left, top, width, height, title, body_lines, accent_color=BLUE_MED):
    """Add a rounded-ish card with accent left bar."""
    # Card background
    card = add_shape_bg(slide, left, top, width, height, GRAY_LIGHT)
    # Accent bar
    add_shape_bg(slide, left, top, Inches(0.08), height, accent_color)
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.5), Inches(0.4),
                 title, font_size=14, bold=True, color=BLUE_DARK)
    # Body
    y = top + Inches(0.55)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.3),
                     line, font_size=11, color=GRAY)
        y += Inches(0.28)

def title_slide(prs, title, subtitle):
    """Create a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, BLUE_DARK)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
                 title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                 subtitle, font_size=20, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    # Bottom line
    add_shape_bg(slide, Inches(4.5), Inches(4.5), Inches(4), Inches(0.05), BLUE_MED)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "Projet SUIVI DE CHANTIERS  •  CAA Guadeloupe", font_size=14, color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)
    return slide

def section_slide(prs, number, title, subtitle=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), BLUE_DARK)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(1.5), Inches(1.6), Inches(1.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE_MED
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    tf.paragraphs[0].space_before = Pt(14)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 title, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                     subtitle, font_size=16, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    return slide

def content_slide(prs, title, subtitle=""):
    """Create a content slide with title bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title_bar(slide)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
                 title, font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
                     subtitle, font_size=13, color=RGBColor(0xBF, 0xDB, 0xFE))
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 1 : TECHNIQUE / ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════

def build_tech_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: Title ──
    title_slide(prs,
        "Suivi de Chantiers — Architecture Technique",
        "Comment l'application est construite, déployée et maintenue")

    # ── SLIDE 2: Sommaire ──
    slide = content_slide(prs, "Sommaire")
    items = [
        "Vue d'ensemble de l'application",
        "C'est quoi ces outils ? (GitHub, Docker, Windsurf, Nginx…)",
        "Stack technique (Frontend + Backend)",
        "Architecture du code source",
        "Base de données (Prisma + SQLite)",
        "Liste complète des dépendances",
        "Infrastructure : VM Debian + Docker",
        "Déploiement : GitHub → Docker Compose",
        "GitHub + Windsurf + Docker : comment ça s'articule",
        "Mettre à jour l'application",
        "Exporter un chat Windsurf pour travailler en local",
    ]
    add_bullet_slide(slide, items, font_size=16, top=Inches(1.8))

    # ── SLIDE 3: Vue d'ensemble ──
    section_slide(prs, 1, "Vue d'ensemble", "Qu'est-ce que l'application fait ?")

    slide = content_slide(prs, "Vue d'ensemble de l'application")
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.5),
                 "Application web de planification et suivi de chantiers pour CAA Guadeloupe", font_size=16, bold=True, color=BLUE_DARK)
    items = [
        "Planification quotidienne : affecter des employés sur des chantiers (AM / PM / Journée)",
        "Vue Gantt : timeline des chantiers avec dates contractuelles / prévisionnelles / réelles",
        "Dashboard chantier : pilotage par projet avec suivi des intervenants",
        "Calendrier personnel : historique individuel d'un employé (absences, affectations)",
        "Liste de plans : gestion des plans liés aux OF et articles",
        "Questions : suivi des questions avec pièces jointes",
        "Création d'OF : partitionnement visuel des articles en ordres de fabrication",
        "Administration : filiales, ateliers, équipes, personnel, chantiers, OF",
        "Synchronisation Fabric : import automatique depuis le Data Warehouse Microsoft Fabric",
    ]
    add_bullet_slide(slide, items, top=Inches(2.4), font_size=13)

    # ══════════════════════════════════════════════════════════════════════════
    # SECTION : C'est quoi ces outils ?
    # ══════════════════════════════════════════════════════════════════════════
    section_slide(prs, 2, "C'est quoi ces outils ?", "Explication simple de chaque technologie utilisée")

    # ── GitHub ──
    slide = content_slide(prs, "GitHub — Le coffre-fort du code", "À quoi ça sert ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.2),
             "En une phrase", [
                 "GitHub est un site web qui stocke le code source du projet en ligne, avec tout l'historique des modifications.",
             ], accent_color=BLUE_MED)
    items = [
        "Versioning : chaque modification (commit) est enregistrée — on peut revenir en arrière à tout moment",
        "Collaboration : plusieurs personnes peuvent travailler sur le même code sans se marcher dessus",
        "Sauvegarde : même si ton PC tombe en panne, le code est en sécurité sur GitHub",
        "Déploiement : la VM de production tire le code depuis GitHub pour mettre à jour l'application",
        "Visibilité : on voit qui a modifié quoi, quand, et pourquoi (message du commit)",
        "",
        "Dans notre projet : le repo est https://github.com/CYBERTESTCAA/PLANNING_PREVISIONNEL",
        "Workflow : tu modifies le code → git add . → git commit → git push → c'est sur GitHub !",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(3.1), spacing=Pt(4))

    # ── Docker ──
    slide = content_slide(prs, "Docker — L'emballage universel", "À quoi ça sert ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.2),
             "En une phrase", [
                 "Docker « emballe » l'application + toutes ses dépendances dans un conteneur — ça marche pareil partout.",
             ], accent_color=GREEN)
    items = [
        "Problème résolu : « ça marche sur mon PC mais pas sur le serveur » → Docker élimine ce problème",
        "Conteneur : c'est comme une mini-machine virtuelle ultra-légère qui contient juste ce qu'il faut",
        "Image : la recette pour créer un conteneur (définie dans le Dockerfile)",
        "Docker Compose : orchestrer plusieurs conteneurs ensemble (ex: frontend + backend)",
        "Isolation : chaque conteneur est indépendant, ils ne se perturbent pas entre eux",
        "",
        "Dans notre projet : 2 conteneurs Docker",
        "  - planning-web : le site web (React compilé + Nginx qui le sert)",
        "  - planning-api : le serveur API (Node.js + Fastify + base de données SQLite)",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(3.1), spacing=Pt(4))

    # ── Windsurf ──
    slide = content_slide(prs, "Windsurf — L'IDE intelligent", "À quoi ça sert ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.2),
             "En une phrase", [
                 "Windsurf est un éditeur de code (comme VS Code) avec une IA intégrée (Cascade) qui aide à écrire du code.",
             ], accent_color=ORANGE)
    items = [
        "IDE : Integrated Development Environment — c'est l'outil où on écrit et modifie le code",
        "Cascade : l'assistant IA intégré à Windsurf qui peut écrire, corriger et expliquer du code",
        "Avantage : on décrit ce qu'on veut en français, Cascade propose le code et l'applique directement",
        "Chat exportable : on peut exporter l'historique de conversation pour qu'un autre développeur comprenne le contexte",
        "",
        "Dans notre projet : l'application a été développée principalement via Windsurf + Cascade",
        "Alternativement : on peut aussi utiliser VS Code (gratuit, sans IA intégrée) ou tout autre éditeur",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(3.1), spacing=Pt(4))

    # ── mRemoteNG + SSH ──
    slide = content_slide(prs, "mRemoteNG + SSH — Accès à distance", "À quoi ça sert ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.2),
             "En une phrase", [
                 "SSH permet de se connecter à un serveur distant en ligne de commande. mRemoteNG est un client qui facilite ça.",
             ], accent_color=BLUE_MED)
    items = [
        "SSH (Secure Shell) : protocole pour se connecter de façon sécurisée à un serveur Linux distant",
        "mRemoteNG : logiciel Windows qui gère les connexions SSH (+ RDP, VNC…) dans des onglets",
        "Pourquoi : la VM Debian (serveur de production) n'a pas d'écran — on y accède uniquement via SSH",
        "Ce qu'on fait en SSH : lancer des commandes Docker, voir les logs, faire des mises à jour manuelles",
        "",
        "Exemple : ssh utilisateur@192.168.13.51 → on est sur la VM → docker compose logs -f api",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(3.1), spacing=Pt(4))

    # ── Nginx ──
    slide = content_slide(prs, "Nginx — Le portier de l'application", "À quoi ça sert ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.2),
             "En une phrase", [
                 "Nginx est un serveur web qui sert les pages du site et redirige les appels API vers le backend.",
             ], accent_color=RED)
    items = [
        "Serveur web : quand un navigateur demande la page, c'est Nginx qui lui envoie les fichiers HTML/CSS/JS",
        "Proxy reverse : quand le navigateur appelle /api/…, Nginx redirige vers le backend Fastify (port 3001)",
        "HTTPS : Nginx gère le certificat TLS pour que la connexion soit chiffrée (cadenas vert)",
        "Sécurité : Nginx ajoute les headers de sécurité (CSP, HSTS, X-Frame-Options…)",
        "Cache : les fichiers statiques (JS/CSS compilés) sont mis en cache 1 an pour accélérer le chargement",
        "",
        "Dans notre projet : Nginx tourne dans le conteneur Docker « planning-web »",
        "Config : docker/nginx.conf (proxy, TLS, headers de sécurité, cache)",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(3.1), spacing=Pt(4))

    # ── Node.js + Fastify ──
    slide = content_slide(prs, "Node.js + Fastify — Le moteur du backend", "À quoi ça sert ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(4.5),
             "Node.js", [
                 "C'est quoi : un runtime JavaScript",
                 "côté serveur (pas dans le navigateur).",
                 "",
                 "Pourquoi : permet d'utiliser le même",
                 "langage (TypeScript/JS) pour le",
                 "frontend ET le backend.",
                 "",
                 "Version : Node.js 20 (LTS)",
             ], accent_color=BLUE_MED)
    add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.5),
             "Fastify", [
                 "C'est quoi : un framework web pour",
                 "Node.js, comme Express mais plus rapide.",
                 "",
                 "Il fait quoi : il écoute les requêtes HTTP",
                 "(GET /plans, POST /questions, etc.)",
                 "et renvoie les données en JSON.",
                 "",
                 "Routes : 20 fichiers dans api/src/routes/",
                 "Chaque fichier = un module (plans, questions…)",
             ], accent_color=GREEN)

    # ── React + Vite + Tailwind ──
    slide = content_slide(prs, "React + Vite + TailwindCSS — Le frontend", "À quoi ça sert ?")
    add_card(slide, Inches(0.3), Inches(1.7), Inches(3.9), Inches(4.5),
             "React", [
                 "C'est quoi : une bibliothèque",
                 "JavaScript pour construire",
                 "des interfaces utilisateur.",
                 "",
                 "Principe : on écrit des",
                 "« composants » réutilisables",
                 "(boutons, modales, grilles…)",
                 "qui se mettent à jour tout seuls",
                 "quand les données changent.",
             ], accent_color=BLUE_MED)
    add_card(slide, Inches(4.5), Inches(1.7), Inches(3.9), Inches(4.5),
             "Vite", [
                 "C'est quoi : un outil de build",
                 "ultra-rapide pour le dev web.",
                 "",
                 "En dev : recharge la page",
                 "instantanément quand on",
                 "modifie le code (Hot Reload).",
                 "",
                 "En prod : compile tout le",
                 "code en fichiers optimisés",
                 "(HTML/JS/CSS minifiés).",
             ], accent_color=GREEN)
    add_card(slide, Inches(8.7), Inches(1.7), Inches(3.9), Inches(4.5),
             "TailwindCSS", [
                 "C'est quoi : un framework CSS",
                 "utility-first.",
                 "",
                 "Au lieu d'écrire du CSS dans",
                 "des fichiers séparés, on met",
                 "des classes directement dans",
                 "le HTML : class=\"bg-blue-500",
                 "text-white p-4 rounded-lg\"",
                 "",
                 "Résultat : développement",
                 "beaucoup plus rapide.",
             ], accent_color=ORANGE)

    # ── Prisma + SQLite ──
    slide = content_slide(prs, "Prisma + SQLite — La base de données", "À quoi ça sert ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(4.5),
             "SQLite", [
                 "C'est quoi : une base de données",
                 "simple, stockée dans UN SEUL fichier.",
                 "",
                 "Avantage : pas besoin d'installer",
                 "un serveur de base de données",
                 "(comme MySQL ou PostgreSQL).",
                 "",
                 "Fichier : planning.sqlite dans",
                 "le volume Docker api-data.",
                 "",
                 "Limite : 1 seul serveur (pas de",
                 "réplication), mais largement",
                 "suffisant pour notre usage.",
             ], accent_color=BLUE_MED)
    add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.5),
             "Prisma (ORM)", [
                 "C'est quoi : un outil qui fait le",
                 "lien entre le code TypeScript",
                 "et la base de données SQLite.",
                 "",
                 "Au lieu d'écrire du SQL :",
                 "  SELECT * FROM plans WHERE ...",
                 "On écrit en TypeScript :",
                 "  prisma.plan.findMany({ where: ... })",
                 "",
                 "Schema : api/prisma/schema.prisma",
                 "définit toutes les tables et relations.",
                 "",
                 "Prisma Studio : interface graphique",
                 "pour explorer la base (prisma studio).",
             ], accent_color=GREEN)

    # ── TypeScript ──
    slide = content_slide(prs, "TypeScript — Du JavaScript mais en mieux", "À quoi ça sert ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.2),
             "En une phrase", [
                 "TypeScript = JavaScript + typage. On déclare les types (string, number, etc.) → les erreurs sont détectées avant l'exécution.",
             ], accent_color=BLUE_MED)
    items = [
        "JavaScript : le langage du web (navigateur). Tous les sites web modernes utilisent du JavaScript",
        "Problème du JS : pas de vérification de type → on découvre les bugs en production",
        "TypeScript : ajoute un système de types → l'éditeur détecte les erreurs pendant qu'on code",
        "Compilation : TypeScript est compilé en JavaScript avant d'être exécuté (le navigateur ne comprend pas TS)",
        "",
        "Dans notre projet : TOUT le code est en TypeScript (frontend + backend)",
        "Fichier tsconfig.json : configure les règles de compilation TypeScript",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(3.1), spacing=Pt(4))

    # ── Azure AD / MSAL / Fabric ──
    slide = content_slide(prs, "Azure AD + Fabric — L'écosystème Microsoft", "À quoi ça sert ?")
    add_card(slide, Inches(0.3), Inches(1.7), Inches(3.9), Inches(4.5),
             "Azure AD (Entra ID)", [
                 "C'est quoi : le système",
                 "d'authentification Microsoft.",
                 "",
                 "Les utilisateurs se connectent",
                 "avec leur compte Microsoft",
                 "professionnel (@caa...).",
                 "",
                 "MSAL : la librairie côté frontend",
                 "qui gère la connexion (popup,",
                 "tokens JWT, etc.).",
             ], accent_color=BLUE_MED)
    add_card(slide, Inches(4.5), Inches(1.7), Inches(3.9), Inches(4.5),
             "Microsoft Fabric", [
                 "C'est quoi : la plateforme",
                 "de données Microsoft (Data",
                 "Warehouse + Analytics).",
                 "",
                 "Dans notre projet : c'est la",
                 "source de vérité pour les",
                 "filiales, ateliers, employés,",
                 "chantiers, OF, articles, etc.",
                 "",
                 "Sync : toutes les heures, l'API",
                 "se connecte à Fabric et importe",
                 "les données nouvelles/modifiées.",
             ], accent_color=ORANGE)
    add_card(slide, Inches(8.7), Inches(1.7), Inches(3.9), Inches(4.5),
             "Connexion Fabric", [
                 "Protocole : SQL (TDS) via",
                 "le package npm « mssql ».",
                 "",
                 "Auth : Service Principal Azure",
                 "(client_id + client_secret).",
                 "",
                 "Endpoint : un serveur SQL",
                 "Fabric exposé par Microsoft.",
                 "",
                 "Config : dans api/.env",
                 "(FABRIC_SQL_ENDPOINT,",
                 "FABRIC_TENANT_ID, etc.)",
             ], accent_color=GREEN)

    # ── Résumé visuel : comment tout s'articule ──
    slide = content_slide(prs, "Résumé — Comment tout s'articule", "La vue d'ensemble de tous les outils")
    schema_lines = [
        "",
        "   DÉVELOPPEUR                              SERVEUR DE PRODUCTION",
        "   ──────────                               ─────────────────────",
        "",
        "   Windsurf (IDE)                           VM Debian 13 (192.168.13.51)",
        "   ├── Code React (frontend)                ├── Docker Compose",
        "   │   ├── TypeScript                       │   ├── planning-web (Nginx)",
        "   │   ├── TailwindCSS                      │   │   ├── Sert le site React compilé",
        "   │   └── Vite (dev server)                │   │   └── Proxy /api → planning-api",
        "   ├── Code Fastify (backend)               │   └── planning-api (Node.js)",
        "   │   ├── Prisma (ORM)                     │       ├── Fastify (API REST)",
        "   │   └── Routes API                       │       ├── Prisma + SQLite",
        "   └── git push ──→ GitHub ──→              │       └── Sync Fabric (cron)",
        "                                            └── mRemoteNG / SSH pour administrer",
        "",
        "   ┌──────────────────────────────────────────────────────────────┐",
        "   │  GitHub = le pont entre le PC dev et le serveur production  │",
        "   │  Docker = garantit que ça tourne pareil partout             │",
        "   │  Nginx  = sert le site web et protège l'accès              │",
        "   └──────────────────────────────────────────────────────────────┘",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(12.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(schema_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.name = "Consolas"
        run.font.color.rgb = BLUE_DARK

    # ── SLIDE: Stack technique ──
    section_slide(prs, 3, "Stack Technique", "Frontend + Backend")

    slide = content_slide(prs, "Stack Technique")
    # Frontend card
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(4.8),
             "FRONTEND (SPA)", [
                 "React 18 + TypeScript",
                 "Vite 5 (dev server + build)",
                 "TailwindCSS 3 (styling utility-first)",
                 "shadcn/ui + Radix UI (composants)",
                 "react-router-dom v6 (navigation)",
                 "FullCalendar (calendrier)",
                 "Recharts (graphiques)",
                 "Framer Motion (animations)",
                 "Lucide React (icônes)",
                 "date-fns (dates)",
                 "jsPDF + jspdf-autotable (export PDF)",
                 "xlsx (export Excel)",
                 "Azure MSAL (authentification Azure AD)",
                 "Zod (validation)",
             ], accent_color=BLUE_MED)
    # Backend card
    add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.8),
             "BACKEND (API REST)", [
                 "Node.js 20 + TypeScript",
                 "Fastify 5 (framework HTTP)",
                 "Prisma 6 (ORM)",
                 "SQLite (base de données)",
                 "Zod (validation requêtes)",
                 "@fastify/cors (CORS)",
                 "@fastify/helmet (sécurité headers)",
                 "@fastify/rate-limit (anti-abus)",
                 "@fastify/multipart (upload fichiers)",
                 "@fastify/static (servir fichiers)",
                 "mssql (connexion Fabric SQL)",
                 "@azure/identity (auth Azure)",
                 "jsonwebtoken + jwks-rsa (JWT)",
                 "node-cron (tâches planifiées)",
             ], accent_color=GREEN)

    # ── SLIDE 5: Architecture code ──
    section_slide(prs, 4, "Architecture du Code Source", "Organisation des fichiers")

    slide = content_slide(prs, "Arborescence du projet")
    code_lines = [
        "shiftwise-planner-main/",
        "├── src/                          ← Code source frontend (React)",
        "│   ├── pages/                    ← Pages de l'application",
        "│   │   ├── PlanningPage.tsx       (grille de planning)",
        "│   │   ├── PlansPage.tsx          (liste des plans)",
        "│   │   ├── QuestionsPage.tsx      (gestion des questions)",
        "│   │   ├── ChantierPage.tsx       (vue chantier)",
        "│   │   ├── GanttPage.tsx          (vue Gantt)",
        "│   │   ├── CreerOFPage.tsx        (création d'OF)",
        "│   │   ├── DashboardChantierPage  (dashboard projet)",
        "│   │   ├── PersonCalendarPage     (calendrier personnel)",
        "│   │   └── admin/                 (pages d'administration)",
        "│   ├── components/               ← Composants réutilisables",
        "│   ├── lib/api.ts                ← Client API typé (fetch wrapper)",
        "│   └── contexts/                 ← State management (PlanningStore)",
        "├── api/                          ← Code source backend (Fastify)",
        "│   ├── src/routes/               ← 20 fichiers de routes API REST",
        "│   ├── src/services/             ← Logique métier (sync Fabric…)",
        "│   ├── prisma/schema.prisma      ← Schéma de la base de données",
        "│   └── prisma/migrations/        ← Historique des migrations",
        "├── docker-compose.yml            ← Orchestration des conteneurs",
        "├── Dockerfile                    ← Build frontend (Nginx)",
        "├── api/Dockerfile                ← Build backend (Node)",
        "├── docker/nginx.conf             ← Config Nginx (proxy, TLS, headers)",
        "└── scripts/                      ← Scripts de déploiement et maintenance",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.name = "Consolas"
        run.font.color.rgb = BLACK if "←" not in line else BLUE_DARK
        if i == 0:
            run.font.bold = True

    # ── SLIDE 6: Routes API ──
    slide = content_slide(prs, "Routes API REST", "20 modules de routes backend")
    routes = [
        "/subsidiaries : CRUD filiales",
        "/workshops : CRUD ateliers",
        "/teams : CRUD équipes",
        "/employees : CRUD personnel",
        "/projects : CRUD chantiers/projets",
        "/manufacturing-orders : CRUD ordres de fabrication",
        "/plans : CRUD plans (création unitaire + bulk)",
        "/assignments : Affectations employé ↔ chantier",
        "/absences : Gestion des absences (CP, RTT, maladie…)",
        "/tasks : Tâches collaborateur",
        "/questions : Questions avec pièces jointes",
        "/attachments : Upload / suppression de fichiers",
        "/planning : Chargement de période planning",
        "/calendar : Données calendrier",
        "/sync : Synchronisation Fabric (trigger + status + history)",
        "/clients : Gestion des clients",
        "/dessinateurs : Gestion des dessinateurs",
        "/affaires : Gestion des affaires",
        "/time-entries : Saisie des temps",
        "/warehouse : Données entrepôt",
    ]
    add_bullet_slide(slide, routes, font_size=12, top=Inches(1.7), spacing=Pt(3))

    # ── SLIDE 7: Base de données ──
    section_slide(prs, 5, "Base de Données", "Prisma ORM + SQLite")

    slide = content_slide(prs, "Modèle de données (simplifié)")
    schema_lines = [
        "Subsidiary (Filiale) ──→ Workshop (Atelier) ──→ Team (Équipe)",
        "                                    │                     │",
        "                              Employee (Employé) ←────────┘",
        "                                    │",
        "                   ┌────────────────┼────────────────┐",
        "                   │                │                │",
        "            Assignment         Absence           Task",
        "         (Affectation)     (CP/RTT/Maladie)   (Tâche)",
        "",
        "Project (Chantier) ──→ ManufacturingOrder (OF) ──→ Article",
        "       │",
        "       ├──→ Plan (Plan de fabrication)",
        "       ├──→ Question ──→ QuestionAttachment (Pièces jointes)",
        "       └──→ Assignment (lien Employee ↔ Project ↔ Date ↔ Slot)",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(schema_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.name = "Consolas"
        run.font.color.rgb = BLUE_DARK if any(x in line for x in ["Subsidiary", "Project", "Employee"]) else BLACK

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
                 "Commandes Prisma : npx prisma db push (dev) │ npx prisma migrate deploy (prod) │ npx prisma studio (GUI)",
                 font_size=12, color=GRAY)

    # ── SLIDE 8: Dépendances ──
    section_slide(prs, 6, "Dépendances Complètes", "Toutes les librairies utilisées")

    slide = content_slide(prs, "Dépendances Frontend", "package.json (racine)")
    frontend_deps = [
        "react 18 + react-dom : Framework UI",
        "react-router-dom 6 : Routage SPA",
        "tailwindcss 3 : CSS utility-first",
        "@radix-ui/* : Composants headless (dialog, dropdown, tabs, tooltip…)",
        "lucide-react : Icônes SVG",
        "date-fns : Manipulation de dates",
        "recharts : Graphiques/diagrammes",
        "framer-motion : Animations",
        "@fullcalendar/* : Calendrier interactif",
        "jspdf + jspdf-autotable : Export PDF",
        "xlsx : Export Excel",
        "zod : Validation de schémas",
        "@azure/msal-browser + @azure/msal-react : Auth Azure AD",
        "@tanstack/react-query : Cache et requêtes asynchrones",
        "vite 5 : Build tool + dev server",
        "typescript 5 : Typage statique",
        "vitest : Tests unitaires",
    ]
    add_bullet_slide(slide, frontend_deps, font_size=12, top=Inches(1.7), spacing=Pt(3))

    slide = content_slide(prs, "Dépendances Backend", "api/package.json")
    backend_deps = [
        "fastify 5 : Framework HTTP performant",
        "@prisma/client 6 : ORM pour SQLite",
        "prisma 6 : CLI et moteur de migration",
        "zod : Validation des requêtes API",
        "@fastify/cors : Cross-Origin Resource Sharing",
        "@fastify/helmet : Security headers",
        "@fastify/rate-limit : Protection anti-abus",
        "@fastify/multipart : Upload de fichiers",
        "@fastify/static : Servir des fichiers statiques",
        "mssql : Connecteur SQL Server (pour Fabric)",
        "@azure/identity : Authentification Azure (service principal)",
        "jsonwebtoken + jwks-rsa : Vérification JWT Azure AD",
        "node-cron : Tâches cron (sync Fabric, backups)",
        "dotenv : Variables d'environnement",
        "tsx : Exécution TypeScript en dev (watch mode)",
    ]
    add_bullet_slide(slide, backend_deps, font_size=12, top=Inches(1.7), spacing=Pt(3))

    # ── SLIDE 9: Infrastructure ──
    section_slide(prs, 7, "Infrastructure", "VM Debian + Docker + Nginx")

    slide = content_slide(prs, "Architecture d'infrastructure")
    infra_lines = [
        "┌─── PC Développeur (192.168.13.50) ────────────────────────┐",
        "│                                                           │",
        "│  Code source  ──→  git push  ──→  GitHub (remote)        │",
        "│  IDE: Windsurf / VS Code                                  │",
        "│  mRemoteNG pour SSH vers la VM                            │",
        "│                                                           │",
        "└───────────────────────────────────────────────────────────┘",
        "                         │",
        "                         ▼  (auto-pull toutes les 5 min)",
        "",
        "┌─── VM Debian 13 (192.168.13.51) ─────────────────────────┐",
        "│                                                           │",
        "│   Docker Compose                                          │",
        "│   ┌──────────────────────┐   ┌───────────────────────┐   │",
        "│   │  planning-web  :80   │   │  planning-api  :3001  │   │",
        "│   │  Nginx + SPA React   │──▶│  Fastify + Prisma     │   │",
        "│   │  Proxy /api → :3001  │   │  SQLite (volume)      │   │",
        "│   │  TLS (self-signed)   │   │  Sync Fabric (cron)   │   │",
        "│   └──────────────────────┘   └───────────────────────┘   │",
        "│                                                           │",
        "│   Accessible sur : https://192.168.13.51 (réseau local)  │",
        "└───────────────────────────────────────────────────────────┘",
    ]
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(infra_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.name = "Consolas"
        run.font.color.rgb = BLUE_DARK

    # ── SLIDE 10: Docker ──
    slide = content_slide(prs, "Docker Compose — 2 conteneurs")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(4.5),
             "planning-web (Frontend)", [
                 "Image : nginx:alpine",
                 "Port : 80 (HTTP) + 443 (HTTPS)",
                 "Contient le build Vite (HTML/JS/CSS)",
                 "Nginx sert les fichiers statiques",
                 "Proxy reverse : /api/* → planning-api:3001",
                 "TLS avec certificat self-signed",
                 "Headers de sécurité (CSP, HSTS, X-Frame…)",
                 "Cache agressif sur /assets/ (immutable)",
             ], accent_color=BLUE_MED)
    add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.5),
             "planning-api (Backend)", [
                 "Image : node:20-slim",
                 "Port : 3001 (interne)",
                 "Fastify + Prisma + SQLite",
                 "Volume Docker : api-data (persistance DB)",
                 "Healthcheck : /health toutes les 30s",
                 "Variables d'env : via api/.env",
                 "Sync Fabric : cron toutes les heures",
                 "Backup : cron quotidien à 2h00",
             ], accent_color=GREEN)

    # ── SLIDE 11: GitHub + Windsurf + Docker ──
    section_slide(prs, 8, "GitHub + Windsurf + Docker", "Comment ça s'articule")

    slide = content_slide(prs, "Workflow de développement")
    items = [
        "1. Développement : ouvrir le projet dans Windsurf (ou VS Code)",
        "2. Lancer en local : npm run dev (frontend :8080) + cd api && npm run dev (backend :3001)",
        "3. Modifier le code : éditer les fichiers, Windsurf AI aide au développement",
        "4. Commiter : git add . → git commit -m \"description\" → git push origin main",
        "5. Déploiement auto : la VM vérifie GitHub toutes les 5 min",
        "6. Si nouveau commit détecté : git pull → docker compose build → restart",
        "7. L'application est mise à jour sur https://192.168.13.51",
    ]
    add_bullet_slide(slide, items, font_size=15, top=Inches(1.8))

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
                 "Repo GitHub : https://github.com/CYBERTESTCAA/PLANNING_PREVISIONNEL",
                 font_size=13, bold=True, color=BLUE_MED)

    # ── SLIDE 12: Mise à jour ──
    section_slide(prs, 9, "Mettre à jour l'application", "Déploiement automatique et manuel")

    slide = content_slide(prs, "Mise à jour — 2 méthodes")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(4),
             "Automatique (toutes les 5 min)", [
                 "1. Le timer systemd vérifie GitHub",
                 "2. Si nouveau commit → git pull",
                 "3. docker compose build --no-cache",
                 "4. docker compose down + up -d",
                 "5. Nettoyage anciennes images Docker",
                 "",
                 "Aucune action nécessaire !",
                 "Il suffit de git push sur main.",
             ], accent_color=GREEN)
    add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4),
             "Manuelle (instantanée)", [
                 "1. SSH vers la VM :",
                 "   ssh utilisateur@192.168.13.51",
                 "2. cd /opt/planning",
                 "3. ./scripts/deploy.sh",
                 "",
                 "Ou rebuild direct :",
                 "   docker compose up -d --build",
             ], accent_color=ORANGE)

    # ── SLIDE 13: Commandes utiles ──
    slide = content_slide(prs, "Commandes utiles sur la VM")
    cmds = [
        "docker compose ps : Voir l'état des conteneurs",
        "docker compose logs -f api : Logs du backend en temps réel",
        "docker compose logs -f web : Logs du frontend/Nginx",
        "docker compose restart : Redémarrer tous les services",
        "docker compose up -d --build : Rebuild + redémarrer",
        "docker compose down : Arrêter tout",
        "./scripts/deploy.sh : Mise à jour manuelle",
        "./scripts/backup.sh : Backup manuel de la base de données",
        "systemctl status planning-update.timer : Vérifier le timer auto-update",
        "cat /var/log/planning-deploy.log : Historique des déploiements",
    ]
    add_bullet_slide(slide, cmds, font_size=13, top=Inches(1.7), spacing=Pt(4))

    # ── SLIDE 14: Exporter Windsurf ──
    section_slide(prs, 10, "Exporter un chat Windsurf", "Pour travailler en local et garder le contexte")

    slide = content_slide(prs, "Exporter un historique Windsurf")
    items = [
        "Pourquoi ? : Pour qu'un autre développeur comprenne ce qui a été fait et puisse continuer",
        "",
        "Comment exporter :",
        "1. Ouvrir Windsurf avec le projet",
        "2. Dans le panneau Cascade (chat AI), cliquer sur les 3 points ··· en haut",
        "3. Sélectionner « Export Chat » ou « Export Conversation »",
        "4. Sauvegarder le fichier .md ou .json généré",
        "5. Le partager avec le nouveau développeur",
        "",
        "Comment l'utiliser :",
        "1. Le nouveau développeur clone le repo GitHub",
        "2. Ouvre le projet dans Windsurf",
        "3. Colle le contenu exporté dans un nouveau chat Cascade",
        "4. Cascade comprend le contexte et peut continuer le travail",
        "",
        "Alternative : le summary du Checkpoint (affiché automatiquement quand on reprend un chat)",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(1.7), spacing=Pt(3))

    # ── SLIDE 15: Backups & Sécurité ──
    slide = content_slide(prs, "Sauvegardes & Sécurité")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(4),
             "Sauvegardes", [
                 "Automatique : tous les jours à 2h00",
                 "Stockage : /var/backups/planning/",
                 "Rétention : 30 dernières sauvegardes",
                 "Format : planning_YYYYMMDD.sqlite.gz",
                 "",
                 "Restauration :",
                 "  docker compose stop api",
                 "  gunzip < backup.gz > /tmp/restore.sqlite",
                 "  docker cp restore planning-api:/data/",
                 "  docker compose start api",
             ], accent_color=BLUE_MED)
    add_card(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(4),
             "Sécurité", [
                 "Authentification : Azure AD (MSAL)",
                 "JWT : vérification via JWKS (Azure)",
                 "HTTPS : TLS 1.2/1.3 (certificat self-signed)",
                 "Headers : CSP, HSTS, X-Frame-Options, etc.",
                 "Rate limiting : @fastify/rate-limit",
                 "CORS : configuré pour le domaine",
                 "Helmet : protection headers côté API",
                 ".env ignoré par git (secrets protégés)",
             ], accent_color=RED)

    # ── SLIDE 16: Fin ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_DARK)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "Merci !", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                 "Questions ?", font_size=24, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, Inches(4.5), Inches(4.8), Inches(4), Inches(0.05), BLUE_MED)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "Projet SUIVI DE CHANTIERS  •  CAA Guadeloupe  •  2025-2026", font_size=14, color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)

    path = os.path.join(OUT_DIR, "..", "docs", "Presentation_Technique.pptx")
    prs.save(path)
    print(f"✅ Présentation technique sauvegardée : {path}")
    return path


# ═══════════════════════════════════════════════════════════════════════════════
#  PRESENTATION 2 : GUIDE UTILISATEUR
# ═══════════════════════════════════════════════════════════════════════════════

def build_user_guide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: Title ──
    title_slide(prs,
        "Suivi de Chantiers — Guide Utilisateur",
        "Comment utiliser toutes les fonctionnalités de l'application")

    # ── SLIDE 2: Sommaire ──
    slide = content_slide(prs, "Sommaire")
    items = [
        "Page d'accueil & navigation",
        "Planning prévisionnel (grille)",
        "Vue Gantt (timeline)",
        "Vue Chantier (pilotage projet)",
        "Calendrier personnel",
        "Liste des Plans",
        "Questions (avec pièces jointes)",
        "Création d'OF (partitionnement)",
        "Administration (filiales, ateliers, personnel, chantiers, OF)",
        "Synchronisation Fabric",
    ]
    add_bullet_slide(slide, items, font_size=16, top=Inches(1.8))

    # ── Page d'accueil ──
    section_slide(prs, 1, "Page d'accueil", "Hub de navigation et état du système")

    slide = content_slide(prs, "Page d'accueil")
    items = [
        "Sélection de la filiale active : filtre toutes les vues de l'application",
        "Accès rapide : cliquer sur les cartes pour naviguer (Planning, Gantt, Chantier, Calendrier)",
        "État de synchronisation : dernière sync Fabric, nombre de fiches mises à jour",
        "Bouton « Synchroniser » : déclenche une synchronisation manuelle avec le Data Warehouse",
        "Accès administration : section en bas de page pour gérer les référentiels",
    ]
    add_bullet_slide(slide, items, font_size=14, top=Inches(1.8))

    # ── Planning ──
    section_slide(prs, 2, "Planning Prévisionnel", "La grille de planification quotidienne")

    slide = content_slide(prs, "Planning — Vue d'ensemble")
    items = [
        "Grille interactive : lignes = employés, colonnes = jours",
        "Créneaux : AM (matin), PM (après-midi), FULL (journée entière)",
        "Zoom : molette ou slider pour ajuster la densité",
        "Filtre : par filiale → atelier → équipe",
        "Période : navigation par semaine/mois avec les flèches",
    ]
    add_bullet_slide(slide, items, font_size=14, top=Inches(1.8))

    slide = content_slide(prs, "Planning — Actions possibles")
    items = [
        "Affecter : cliquer sur une cellule vide → sélectionner un chantier → choisir le créneau",
        "Modifier : cliquer sur une affectation existante → changer chantier ou créneau",
        "Supprimer : clic droit → Supprimer, ou sélectionner + touche Suppr",
        "Multi-affectation : sélectionner plusieurs cellules (Ctrl+clic ou sélection en rectangle)",
        "Absences : marquer un jour comme CP, RTT, Maladie, Formation, Autre",
        "Copier/Coller : dupliquer une affectation sur d'autres jours",
        "Undo/Redo : Ctrl+Z / Ctrl+Y pour annuler ou rétablir",
        "Export : bouton Export → PDF ou Excel pour impression",
        "Barre d'actions groupées : actions en masse sur la sélection courante",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(1.8), spacing=Pt(4))

    # ── Gantt ──
    section_slide(prs, 3, "Vue Gantt", "Timeline des chantiers")

    slide = content_slide(prs, "Vue Gantt")
    items = [
        "Timeline horizontale : chaque barre = un chantier",
        "3 couches temporelles par chantier :",
        "  - Contractuel : dates prévues au contrat (bleu)",
        "  - Prévisionnel : dates estimées réelles (orange)",
        "  - Réel : dates effectives de début/fin (vert)",
        "Statut : couleur de la barre selon l'état (en cours, terminé, en retard)",
        "Navigation : scroll horizontal pour parcourir les mois",
        "Clic sur un chantier : ouvre le détail du projet",
    ]
    add_bullet_slide(slide, items, font_size=14, top=Inches(1.8))

    # ── Chantier ──
    section_slide(prs, 4, "Vue Chantier", "Pilotage par projet")

    slide = content_slide(prs, "Vue Chantier")
    items = [
        "Sélection du chantier : dropdown avec recherche par code ou nom",
        "Dashboard : vue synthétique du chantier sélectionné",
        "Intervenants : liste des employés affectés au chantier (AM/PM)",
        "OF liés : ordres de fabrication associés au chantier",
        "Progression : barre d'avancement du projet",
        "Historique : timeline des affectations passées et futures",
    ]
    add_bullet_slide(slide, items, font_size=14, top=Inches(1.8))

    # ── Calendrier personnel ──
    section_slide(prs, 5, "Calendrier Personnel", "Historique individuel")

    slide = content_slide(prs, "Calendrier Personnel")
    items = [
        "Sélection : Filiale → Atelier → Équipe → Personne",
        "Vues : Mois / Semaine / Jour / Année (FullCalendar)",
        "Événements : affectations (couleur du chantier) + absences (icône + type)",
        "Statistiques rapides : nombre d'affectations, absences, chantiers distincts",
        "Navigation : cliquer sur un jour pour voir le détail",
    ]
    add_bullet_slide(slide, items, font_size=14, top=Inches(1.8))

    # ── Plans ──
    section_slide(prs, 6, "Liste des Plans", "Gestion des plans de fabrication")

    slide = content_slide(prs, "Liste des Plans")
    items = [
        "Sélection du chantier : filtrer les plans par projet",
        "Vue panneau gauche : liste des OF avec leurs articles",
        "Badge plan : chaque article affiche le numéro de plan s'il en a un",
        "Articles désactivés : les articles déjà associés à un plan sont grisés",
        "Créer un plan : bouton « Créer depuis articles » ou « + Nouveau plan »",
    ]
    add_bullet_slide(slide, items, font_size=14, top=Inches(1.8))

    slide = content_slide(prs, "Plans — Création d'un plan")
    items = [
        "Sélectionner les articles à inclure dans le plan (cocher les cases)",
        "Restriction : tous les articles doivent appartenir au même OF",
        "Si articles de différents OFs sélectionnés → message d'erreur rouge",
        "Plan sans article : cocher « Créer un plan sans article (cadeau, hors vente…) »",
        "Numéro HK : rempli automatiquement ou saisie manuelle",
        "Numéro de plan : rempli automatiquement (dernier + 1) ou saisie manuelle",
        "Valider : le plan est créé et les articles sont marqués comme « planifiés »",
    ]
    add_bullet_slide(slide, items, font_size=14, top=Inches(1.8))

    # ── Questions ──
    section_slide(prs, 7, "Questions", "Suivi des questions avec pièces jointes")

    slide = content_slide(prs, "Questions")
    items = [
        "Sélection du chantier : filtrer les questions par projet",
        "3 colonnes Kanban : Non traité / En cours / Terminé",
        "Créer une question : bouton « + Nouvelle question »",
        "  → Remplir : désignation, zone, auteur, destinataire, texte de la question",
        "Répondre : cliquer « Répondre » sous la question → écrire la réponse",
        "  → Valider (marque comme terminé) ou « En cours » (garde en progression)",
        "Dates affichées : date de la question + date de réponse (si répondue)",
        "Changer l'état : dropdown d'avancement sur chaque carte",
        "Supprimer : icône poubelle rouge",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(1.8), spacing=Pt(4))

    slide = content_slide(prs, "Questions — Pièces jointes")
    items = [
        "Joindre un fichier : cliquer le bouton « Joindre » sous chaque question",
        "Types acceptés : images (JPG, PNG, GIF…), PDF, Word, Excel, texte",
        "Upload multiple : sélectionner plusieurs fichiers en une fois",
        "Aperçu : les images sont affichées en miniature (64×64px)",
        "Fichiers non-image : affichés avec icône fichier + nom",
        "Ouvrir : cliquer sur la miniature ou le nom pour ouvrir dans un nouvel onglet",
        "Supprimer : survoler la pièce jointe → bouton ✕ rouge en haut à droite",
    ]
    add_bullet_slide(slide, items, font_size=14, top=Inches(1.8))

    # ── Création OF ──
    section_slide(prs, 8, "Création d'OF", "Partitionnement visuel des articles")

    slide = content_slide(prs, "Création d'OF — Fonctionnement")
    items = [
        "Sélection du chantier : dropdown en haut de page",
        "OF existants (synchronisés) : panneau dépliable montrant les OF déjà créés",
        "  → Chaque OF affiche : code, nombre d'articles, nombre de plans",
        "  → Cliquer pour déplier et voir les articles détaillés",
        "  → Bouton « Import » : importe les articles de l'OF dans le pool de travail",
        "Articles avec plan : affichés en semi-transparent avec le numéro de plan",
        "Pool d'articles (gauche) : liste des éléments à organiser, drag & drop possible",
        "Groupes OF (droite) : créer des groupes en saisissant un code + désignation",
        "Glisser-déposer : déplacer les articles du pool vers un groupe OF",
        "Sauvegarde auto : tout est sauvegardé en localStorage automatiquement",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(1.8), spacing=Pt(4))

    # ── Administration ──
    section_slide(prs, 9, "Administration", "Gestion des données de référence")

    slide = content_slide(prs, "Modules d'administration")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(3.8), Inches(4.5),
             "Filiales & Ateliers", [
                 "Créer / modifier / supprimer",
                 "des filiales (ex: CAA Guadeloupe)",
                 "et leurs ateliers.",
                 "",
                 "Vue arborescente :",
                 "Filiale → Ateliers",
             ], accent_color=BLUE_MED)
    add_card(slide, Inches(4.6), Inches(1.7), Inches(3.8), Inches(4.5),
             "Personnel", [
                 "Liste de tous les employés",
                 "Filtre par filiale/atelier",
                 "Activer / désactiver",
                 "",
                 "Pool d'employés :",
                 "non affectés à un atelier",
                 "→ affecter manuellement",
             ], accent_color=GREEN)
    add_card(slide, Inches(8.8), Inches(1.7), Inches(3.8), Inches(4.5),
             "Chantiers & OF", [
                 "Chantiers : code, libellé,",
                 "couleur, statut, progression",
                 "Dates contractuelles/prévi.",
                 "",
                 "Ordres de fabrication :",
                 "liés aux chantiers,",
                 "avec articles associés",
             ], accent_color=ORANGE)

    # ── Sync Fabric ──
    section_slide(prs, 10, "Synchronisation Fabric", "Import automatique depuis le Data Warehouse")

    slide = content_slide(prs, "Synchronisation Fabric")
    items = [
        "Source : Microsoft Fabric Data Warehouse (SQL endpoint)",
        "Données synchronisées : filiales, ateliers, employés, chantiers, OF, articles, clients, affaires",
        "Fréquence : automatique toutes les heures (configurable via FABRIC_SYNC_CRON)",
        "Déclenchement manuel : page d'accueil → bouton « Synchroniser »",
        "Page admin Sync : vue détaillée de l'historique des synchronisations",
        "  → Compteurs : nombre de fiches créées / mises à jour par table",
        "  → Erreurs : affichées en rouge avec le message d'erreur",
        "  → Durée : temps d'exécution de chaque sync",
        "Principe : la sync fait un upsert (création si inexistant, mise à jour sinon)",
        "Les données locales (planning, absences, questions…) ne sont PAS écrasées",
    ]
    add_bullet_slide(slide, items, font_size=13, top=Inches(1.8), spacing=Pt(4))

    # ── SLIDE FIN ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_DARK)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "Bonne utilisation !", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                 "N'hésitez pas à poser des questions", font_size=24, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, Inches(4.5), Inches(4.8), Inches(4), Inches(0.05), BLUE_MED)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "Projet SUIVI DE CHANTIERS  •  CAA Guadeloupe  •  https://192.168.13.51",
                 font_size=14, color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)

    path = os.path.join(OUT_DIR, "..", "docs", "Guide_Utilisateur.pptx")
    prs.save(path)
    print(f"✅ Guide utilisateur sauvegardé : {path}")
    return path


# ─── Main ────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    os.makedirs(os.path.join(OUT_DIR, "..", "docs"), exist_ok=True)
    build_tech_presentation()
    build_user_guide()
    print("\n🎉 Les deux présentations ont été générées dans le dossier docs/")
