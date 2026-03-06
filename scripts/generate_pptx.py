"""
Génère une présentation PowerPoint pour le service IT.
Usage: python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Couleurs ──────────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1E, 0x1E, 0x2E)
BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
DARK_BLUE = RGBColor(0x1E, 0x40, 0x8A)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
GREEN = RGBColor(0x05, 0x96, 0x69)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
RED = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_card(slide, left, top, width, height, title, items, accent_color=BLUE):
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, WHITE)
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    card.line.width = Pt(1)
    # Accent bar
    add_rect(slide, left, top, Inches(0.06), height, accent_color)
    # Title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, font_size=14, bold=True, color=accent_color)
    # Items
    add_bullet_list(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7),
                    items, font_size=11, color=GRAY)

def slide_header(slide, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), BLUE)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.55),
                 title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.4),
                     subtitle, font_size=14, color=LIGHT_BLUE)

def slide_number(slide, num, total):
    add_text_box(slide, Inches(12), Inches(7.05), Inches(1), Inches(0.35),
                 f"{num}/{total}", font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 10

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PAGE DE TITRE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, BLUE)
add_rect(slide, Inches(0), Inches(5.5), prs.slide_width, Inches(2), DARK_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Planning Prévisionnel", font_size=44, bold=True, color=WHITE)
add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.6),
             "Application web de gestion des affectations — Groupe CAA", font_size=22, color=LIGHT_BLUE)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
             "Présentation technique pour le service IT", font_size=16, color=LIGHT_BLUE)

add_text_box(slide, Inches(1), Inches(5.8), Inches(5), Inches(0.4),
             "Niel POUPELIN — Service IT", font_size=14, color=LIGHT_BLUE)
add_text_box(slide, Inches(1), Inches(6.2), Inches(5), Inches(0.4),
             "Mars 2026", font_size=14, color=LIGHT_BLUE)
slide_number(slide, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Sommaire")

sommaire = [
    "1.  Vue d'ensemble de l'application",
    "2.  Architecture technique (Frontend / Backend / BDD)",
    "3.  Dépendances — Frontend",
    "4.  Dépendances — Backend & Infrastructure",
    "5.  Synchronisation Microsoft Fabric",
    "6.  Authentification Azure AD (Entra ID)",
    "7.  Environnement de développement & déploiement",
    "8.  Workflow : GitHub → Windsurf → Docker",
    "9.  Procédure de mise à jour",
    "10. Reprendre le travail / Exporter un chat Windsurf",
]
add_bullet_list(slide, Inches(1), Inches(1.6), Inches(10), Inches(5),
                sommaire, font_size=18, color=BLACK, spacing=Pt(14))
slide_number(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VUE D'ENSEMBLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_header(slide, "Vue d'ensemble", "Qu'est-ce que le Planning Prévisionnel ?")

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
             "Application web interne permettant de planifier les affectations du personnel sur les chantiers, "
             "synchronisée automatiquement avec les données ERP via Microsoft Fabric.",
             font_size=15, color=GRAY)

add_card(slide, Inches(0.8), Inches(2.5), Inches(3.6), Inches(2.3),
         "Fonctionnalités principales",
         ["Planning hebdo par atelier", "Vue Gantt des chantiers",
          "Gestion filiales / ateliers / personnel", "Export PDF des plannings",
          "Absences et demandes de modification"], BLUE)

add_card(slide, Inches(4.8), Inches(2.5), Inches(3.6), Inches(2.3),
         "Données synchronisées",
         ["5 filiales (ST BARTH, GUADELOUPE, SUISSE…)", "~1 000 salariés",
          "~4 200 chantiers / commandes", "Ordres de fabrication",
          "Heures réelles (ERP TimeEntries)"], GREEN)

add_card(slide, Inches(8.8), Inches(2.5), Inches(3.6), Inches(2.3),
         "Accès & Sécurité",
         ["Authentification Azure AD (Entra ID)", "Groupes AD → rôle Admin",
          "HTTPS auto-signé sur VM", "API protégée par JWT",
          "Données en SQLite (volume Docker)"], ORANGE)

add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
             "URL : https://192.168.13.51   ·   VM : SRV-WEB-01 (Debian)   ·   Sync auto toutes les heures",
             font_size=13, bold=True, color=DARK_BLUE)
slide_number(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE TECHNIQUE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Architecture technique", "3 couches : Frontend → API → Base de données")

# Frontend box
add_card(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(3.2),
         "Frontend  (Container: planning-web)",
         ["React 18 + TypeScript + Vite",
          "TailwindCSS + shadcn/ui (Radix)",
          "Framer Motion (animations)",
          "MSAL.js pour Azure AD login",
          "Servi par Nginx (port 80/443)",
          "Build statique dans Docker"], BLUE)

# Arrow
add_text_box(slide, Inches(4.4), Inches(2.8), Inches(0.6), Inches(0.5),
             "→", font_size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# API box
add_card(slide, Inches(4.9), Inches(1.6), Inches(3.8), Inches(3.2),
         "API Backend  (Container: planning-api)",
         ["Fastify 5 + TypeScript",
          "Prisma ORM (migrations auto)",
          "JWT validation (Azure AD keys)",
          "CRUD : filiales, ateliers, employés…",
          "Sync Fabric via node-cron",
          "Port 3001 (localhost only)"], GREEN)

# Arrow
add_text_box(slide, Inches(8.8), Inches(2.8), Inches(0.6), Inches(0.5),
             "→", font_size=32, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# DB box
add_card(slide, Inches(9.3), Inches(1.6), Inches(3.5), Inches(3.2),
         "Base de données",
         ["SQLite en mode WAL",
          "Volume Docker persistant",
          "Schéma Prisma (20+ modèles)",
          "Filiales, ateliers, employés",
          "Clients, affaires, chantiers",
          "Heures, calendrier, jours fériés"], ORANGE)

# Docker compose note
add_rounded_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), LIGHT_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.35),
             "Docker Compose — Orchestration", font_size=14, bold=True, color=DARK_BLUE)
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.9),
             "• docker-compose.yml définit 2 services : api + web\n"
             "• Volume api-data pour la persistence SQLite\n"
             "• Volume /opt/planning/certs pour le certificat SSL\n"
             "• Nginx reverse proxy : /api/* → planning-api:3001, tout le reste → fichiers statiques",
             font_size=12, color=GRAY)
slide_number(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DÉPENDANCES FRONTEND
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Dépendances — Frontend", "package.json (racine du projet)")

add_card(slide, Inches(0.5), Inches(1.6), Inches(4), Inches(2.8),
         "Framework & UI",
         ["react ^18.3  /  react-dom ^18.3",
          "react-router-dom ^6.30",
          "tailwindcss ^3.4  /  postcss",
          "class-variance-authority (CVA)",
          "lucide-react ^0.462 (icônes)",
          "framer-motion ^12.28 (animations)",
          "sonner ^1.7 (toasts)"], BLUE)

add_card(slide, Inches(4.8), Inches(1.6), Inches(4), Inches(2.8),
         "Composants & Données",
         ["@radix-ui/* (15 packages shadcn)",
          "@tanstack/react-query ^5.83",
          "@fullcalendar/* (6 packages)",
          "recharts ^2.15 (graphiques)",
          "date-fns ^3.6 (dates)",
          "zod ^3.25 (validation)",
          "xlsx ^0.18 / jspdf ^4.0 (export)"], GREEN)

add_card(slide, Inches(9.1), Inches(1.6), Inches(3.7), Inches(2.8),
         "Auth & Dev Tools",
         ["@azure/msal-browser ^3.30",
          "@azure/msal-react ^2.2",
          "vite ^5.4 (bundler)",
          "typescript ^5.8",
          "vitest ^3.2 (tests)",
          "eslint ^9.32",
          "autoprefixer ^10.4"], ORANGE)

add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(0.4),
             "Total : ~40 dépendances directes + Radix UI ecosystem (shadcn/ui)",
             font_size=13, bold=True, color=DARK_BLUE)
slide_number(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DÉPENDANCES BACKEND & INFRA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Dépendances — Backend & Infrastructure", "api/package.json + docker-compose.yml")

add_card(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(3.5),
         "Backend (api/package.json)",
         ["fastify ^5.2 (serveur HTTP)",
          "@prisma/client ^6.1 (ORM)",
          "prisma ^6.1 (CLI migrations)",
          "@azure/identity ^4.13 (OAuth)",
          "mssql ^12.2 (Fabric SQL)",
          "jsonwebtoken ^9.0 (JWT)",
          "jwks-rsa ^4.0 (Azure AD keys)",
          "node-cron ^4.2 (planificateur)",
          "dotenv ^16.4  /  zod ^3.25"], BLUE)

add_card(slide, Inches(4.7), Inches(1.6), Inches(3.8), Inches(3.5),
         "Infrastructure VM (SRV-WEB-01)",
         ["Debian Linux (VM interne)",
          "Docker Engine + Docker Compose",
          "Nginx (dans container web)",
          "OpenSSL (certificat auto-signé)",
          "UFW firewall (ports 80, 443, 1433)",
          "Git (pull depuis GitHub)",
          "Node.js 22 (dans containers)",
          "SQLite 3 (via Prisma, dans volume)"], GREEN)

add_card(slide, Inches(8.9), Inches(1.6), Inches(3.8), Inches(3.5),
         "Services externes",
         ["Microsoft Fabric Warehouse",
          "  └ WH_OR_GROUPE_CAA (SQL endpoint)",
          "Azure AD / Entra ID",
          "  └ Tenant: 5598…133e",
          "  └ App Registration: c7ab…160c",
          "  └ Service Principal (sync)",
          "  └ Groupe: GR_CHEFS_ATELIERS",
          "GitHub (CYBERTESTCAA/…)"], ORANGE)

add_rounded_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), LIGHT_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1),
             "Ports réseau ouverts :  80 (HTTP→HTTPS redirect)  ·  443 (HTTPS frontend+API)  ·  1433 (sortant vers Fabric SQL)\n"
             "Variables d'environnement sensibles dans api/.env sur la VM (pas commitées sur GitHub)",
             font_size=12, color=GRAY)
slide_number(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SYNC FABRIC
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_header(slide, "Synchronisation Microsoft Fabric", "Données ERP → SQLite via Service Principal")

steps = [
    ("1", "Cron déclenche", "Toutes les heures\n(ou manuellement)"),
    ("2", "OAuth Token", "Service Principal →\ndatabase.windows.net"),
    ("3", "Requêtes SQL", "9 étapes vers\nFabric Warehouse"),
    ("4", "Import SQLite", "Prisma upsert /\ncreateMany batch"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 3.2)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(1.7), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, x + Inches(0.7), Inches(1.7), Inches(2.2), Inches(0.35),
                 title, font_size=14, bold=True, color=DARK_BLUE)
    add_text_box(slide, x + Inches(0.7), Inches(2.1), Inches(2.2), Inches(0.8),
                 desc, font_size=12, color=GRAY)

# Data table
add_text_box(slide, Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
             "Les 9 étapes de synchronisation :", font_size=14, bold=True, color=BLACK)

sync_steps = [
    "Étape 1 — Entreprises (filiales) : DIM_Entreprise → Subsidiary",
    "Étape 2 — Clients : DIM_Client → Client",
    "Étape 3 — Affaires : DIM_Affaire → Affaire",
    "Étape 4 — Salariés : DIM_Salarie → Employee (enrichi : qualification, service, intérimaire…)",
    "Étape 5 — Commandes : FAIT_Commande → Project (chantier, montants agrégés, noms résolus)",
    "Étape 6 — Ordres de fabrication : FAIT_OF → ManufacturingOrder",
    "Étape 7 — Heures réelles : FAIT_HeureReel → TimeEntry (batch insert)",
    "Étape 8 — Calendrier : DIM_Calendrier → CalendarDay (jours fériés inclus)",
    "Étape 9 — Jours par filiale : DIM_JoursParFiliale → SubsidiarySchedule",
]
add_bullet_list(slide, Inches(0.8), Inches(3.65), Inches(11.5), Inches(3.5),
                sync_steps, font_size=11, color=GRAY, spacing=Pt(4))
slide_number(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — AUTHENTIFICATION AZURE AD
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Authentification Azure AD (Entra ID)", "Login utilisateur + autorisation admin")

# Flow
flow_items = [
    ("Utilisateur", "Ouvre https://192.168.13.51\nClique « Se connecter »", BLUE),
    ("MSAL Redirect", "Redirigé vers login.microsoft.com\nSaisit ses identifiants AD", DARK_BLUE),
    ("Token JWT", "Azure renvoie un ID Token\nContient : nom, email, groups[]", GREEN),
    ("Vérification API", "Backend valide le JWT\nVérifie signature + groupe admin", ORANGE),
]
for i, (title, desc, color) in enumerate(flow_items):
    x = Inches(0.5 + i * 3.2)
    add_rounded_rect(slide, x, Inches(1.6), Inches(2.8), Inches(1.8), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), Inches(1.7), Inches(2.5), Inches(0.4),
                 title, font_size=13, bold=True, color=color)
    add_text_box(slide, x + Inches(0.15), Inches(2.1), Inches(2.5), Inches(1.2),
                 desc, font_size=11, color=GRAY)
    if i < 3:
        add_text_box(slide, x + Inches(2.85), Inches(2.2), Inches(0.4), Inches(0.4),
                     "→", font_size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.4),
             "Configuration Azure AD :", font_size=14, bold=True, color=BLACK)

config_items = [
    "Tenant ID : 55981225-247f-4ff6-8678-e1efb27d133e",
    "Client ID (App Registration) : c7abec75-508c-4495-9043-851280c3160c",
    "Groupe Admin : GR_CHEFS_ATELIERS (bd3a6f77-0a07-4ade-a2e5-d26cd8310297)",
    "Redirect URI : https://192.168.13.51  (configuré dans l'App Registration Azure)",
    "Flow : MSAL loginRedirect → ID Token → API valide via JWKS (keys Microsoft)",
    "Admin = membre du groupe GR_CHEFS_ATELIERS → accès section Administration",
]
add_bullet_list(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.5),
                config_items, font_size=12, color=GRAY, spacing=Pt(4))
slide_number(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ENVIRONNEMENT & WORKFLOW
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_header(slide, "Workflow : GitHub + Windsurf + Docker", "Du code au déploiement")

# Dev environment
add_card(slide, Inches(0.5), Inches(1.6), Inches(4), Inches(2.5),
         "Développement (PC local)",
         ["IDE : Windsurf (fork VS Code + AI)",
          "Cascade AI assiste le codage",
          "Git : commit + push vers GitHub",
          "Repo : CYBERTESTCAA/PLANNING_PREVISIONNEL",
          "Branche : main"], BLUE)

# Arrow
add_text_box(slide, Inches(4.6), Inches(2.6), Inches(0.6), Inches(0.5),
             "→", font_size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# GitHub
add_card(slide, Inches(5.1), Inches(1.6), Inches(3), Inches(2.5),
         "GitHub",
         ["Repository privé",
          "Code source versionné",
          "Historique des commits",
          "Pas de CI/CD automatique",
          "(déploiement manuel)"], DARK_BLUE)

# Arrow
add_text_box(slide, Inches(8.2), Inches(2.6), Inches(0.6), Inches(0.5),
             "→", font_size=32, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# VM
add_card(slide, Inches(8.7), Inches(1.6), Inches(4), Inches(2.5),
         "VM SRV-WEB-01 (Debian)",
         ["Accès : MRemoteNG + SSH",
          "cd /opt/planning && git pull",
          "docker compose build",
          "docker compose up -d",
          "Données persistantes (volume)"], GREEN)

# Update procedure
add_rounded_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), WHITE)
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.4),
             "Procédure de mise à jour complète :", font_size=14, bold=True, color=DARK_BLUE)

update_steps = [
    "1. Sur le PC : modifier le code dans Windsurf → tester localement (npm run dev)",
    "2. Commit & push :   git add -A && git commit -m \"description\" && git push",
    "3. Connexion SSH à SRV-WEB-01 via MRemoteNG (192.168.13.51)",
    "4. Pull le code :   cd /opt/planning && git pull",
    "5. Rebuild les containers :   docker compose build",
    "6. Relancer :   docker compose up -d",
    "7. Vérifier les logs :   docker compose logs -f --tail 50",
    "8. Les données SQLite sont conservées (volume Docker persistant)",
]
add_bullet_list(slide, Inches(0.8), Inches(5.05), Inches(11.5), Inches(1.7),
                update_steps, font_size=11, color=GRAY, spacing=Pt(2))
slide_number(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — EXPORTER CHAT WINDSURF + CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Reprendre le travail / Exporter un chat Windsurf")

add_card(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.8),
         "Exporter un chat Windsurf (Cascade)",
         ["Dans Windsurf, ouvrir le panneau Cascade (chat AI)",
          "Cliquer sur l'icône ⋯ (3 points) en haut du chat",
          "Sélectionner « Export Conversation »",
          "Sauvegarde un fichier .md avec tout l'historique",
          "Contient : requêtes, réponses, fichiers modifiés",
          "→ Le nouveau développeur peut lire le contexte complet"], BLUE)

add_card(slide, Inches(6.7), Inches(1.6), Inches(6), Inches(2.8),
         "Reprendre le projet sur un autre poste",
         ["1. Cloner le repo : git clone https://github.com/CYBERTESTCAA/…",
          "2. npm install (installe toutes les dépendances)",
          "3. Copier .env depuis la VM ou un collègue",
          "4. npm run dev → lance le frontend (localhost:8080)",
          "5. npm run dev:api → lance l'API (localhost:3001)",
          "6. Lire le chat Windsurf exporté pour le contexte"], GREEN)

add_card(slide, Inches(0.5), Inches(4.8), Inches(12.2), Inches(1.8),
         "Fichiers importants à connaître",
         ["docker-compose.yml — orchestration des 2 containers (api + web)",
          "api/.env — variables sensibles (Fabric credentials, Azure AD IDs) — NE PAS COMMITER",
          "api/src/services/fabricSync.ts — logique de synchronisation Fabric → SQLite",
          "api/src/middleware/auth.ts — validation JWT Azure AD + vérification groupe admin",
          "src/contexts/PlanningStoreContext.tsx — store central React (charge les données API au démarrage)",
          "src/contexts/AuthContext.tsx — gestion MSAL login/logout + token provider"], ORANGE)
slide_number(slide, 10, TOTAL_SLIDES)

# ── Sauvegarde ────────────────────────────────────────────────────────────────
output = "Presentation_Planning_Previsionnel_IT.pptx"
prs.save(output)
print(f"✅ Présentation générée : {output}")
