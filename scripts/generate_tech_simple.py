#!/usr/bin/env python3
"""
Generate a user-friendly technical presentation for non-developers.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BLUE_DARK  = RGBColor(0x1E, 0x3A, 0x5F)
BLUE_MED   = RGBColor(0x25, 0x63, 0xEB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1F, 0x1F, 0x1F)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
GREEN      = RGBColor(0x05, 0x96, 0x69)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
RED        = RGBColor(0xDC, 0x26, 0x26)

def add_bg(slide, color=WHITE):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = alignment
    run = p.add_run(); run.text = text; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.color.rgb = color; run.font.name = "Calibri"
    return txBox

def add_bullets(slide, items, left=Inches(0.8), top=Inches(2.0), width=Inches(11.5), font_size=15, color=BLACK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        if " : " in item:
            parts = item.split(" : ", 1)
            r = p.add_run(); r.text = "•  " + parts[0] + " : "; r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
            r2 = p.add_run(); r2.text = parts[1]; r2.font.size = Pt(font_size); r2.font.color.rgb = color; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = "•  " + item; r.font.size = Pt(font_size); r.font.color.rgb = color; r.font.name = "Calibri"

def add_card(slide, left, top, width, height, title, lines, accent=BLUE_MED):
    add_shape_bg(slide, left, top, width, height, GRAY_LIGHT)
    add_shape_bg(slide, left, top, Inches(0.08), height, accent)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.5), Inches(0.4), title, font_size=14, bold=True, color=BLUE_DARK)
    y = top + Inches(0.55)
    for line in lines:
        add_text_box(slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.3), line, font_size=11, color=GRAY)
        y += Inches(0.28)

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_DARK)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2), title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8), subtitle, font_size=20, color=RGBColor(0xBF,0xDB,0xFE), alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, Inches(4.5), Inches(4.5), Inches(4), Inches(0.05), BLUE_MED)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "Projet SUIVI DE CHANTIERS  •  CAA Guadeloupe", font_size=14, color=RGBColor(0x93,0xC5,0xFD), alignment=PP_ALIGN.CENTER)

def section_slide(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), BLUE_DARK)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(1.5), Inches(1.6), Inches(1.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = BLUE_MED; circle.line.fill.background()
    tf = circle.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(number); r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    tf.paragraphs[0].space_before = Pt(14)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1), title, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6), subtitle, font_size=16, color=RGBColor(0xBF,0xDB,0xFE), alignment=PP_ALIGN.CENTER)

def content_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5), title, font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4), subtitle, font_size=13, color=RGBColor(0xBF,0xDB,0xFE))
    return slide

def mono_block(slide, lines, top=Inches(1.5)):
    txBox = slide.shapes.add_textbox(Inches(0.3), top, Inches(12.5), Inches(5.5))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run(); r.text = line; r.font.size = Pt(11); r.font.name = "Consolas"; r.font.color.rgb = BLUE_DARK


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title ──
    title_slide(prs, "Suivi de Chantiers — Comment ca marche ?", "Explique simplement, pour tout le monde")

    # ── Sommaire ──
    slide = content_slide(prs, "Ce qu'on va voir ensemble")
    add_bullets(slide, [
        "L'application en bref : a quoi elle sert",
        "Comment ca fonctionne : l'analogie du restaurant",
        "Les outils : a quoi sert chaque brique",
        "Ou vit l'application : le serveur",
        "D'ou viennent les donnees : la synchronisation HK",
        "Comment on met a jour l'application",
        "La securite : qui peut acceder",
        "Pour aller plus loin : si vous devez intervenir",
    ], font_size=18, top=Inches(1.8))

    # ═══════════ SECTION 1 : L'app en bref ═══════════
    section_slide(prs, 1, "L'application en bref", "Ce qu'elle fait, pour qui, et pourquoi")

    slide = content_slide(prs, "Qu'est-ce que l'application « Suivi de Chantiers » ?")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.0), "En une phrase", [
        "C'est un site web interne qui permet de planifier qui travaille sur quel chantier, et de suivre l'avancement.",
    ], accent=BLUE_MED)
    add_bullets(slide, [
        "Planning : affecter les employes sur les chantiers, matin ou apres-midi",
        "Vue Gantt : voir les dates de debut/fin de chaque chantier sur un calendrier",
        "Suivi chantier : savoir qui intervient et ou on en est",
        "Calendrier personnel : voir le planning d'un employe (presences, absences)",
        "Plans : gerer les plans de fabrication lies aux commandes",
        "Questions : poser des questions sur un chantier et suivre les reponses",
        "Creation d'OF : regrouper les articles en ordres de fabrication",
    ], top=Inches(3.0), font_size=13, spacing=Pt(5))

    slide = content_slide(prs, "Qui utilise l'application ?")
    add_card(slide, Inches(0.3), Inches(1.7), Inches(3.9), Inches(3.5), "Les chefs d'atelier", [
        "Ils planifient les employes", "sur les chantiers au quotidien.", "",
        "Ils voient les absences, les", "disponibilites, les chantiers", "en cours et a venir.",
    ], accent=BLUE_MED)
    add_card(slide, Inches(4.5), Inches(1.7), Inches(3.9), Inches(3.5), "Les responsables", [
        "Ils suivent l'avancement", "global des chantiers.", "",
        "Ils consultent la vue Gantt,", "les dashboards, et les", "questions en cours.",
    ], accent=GREEN)
    add_card(slide, Inches(8.7), Inches(1.7), Inches(3.9), Inches(3.5), "Le bureau d'etudes", [
        "Ils gerent les plans de", "fabrication et les ordres", "de fabrication (OF).", "",
        "Ils posent et repondent", "aux questions techniques.",
    ], accent=ORANGE)

    # ═══════════ SECTION 2 : Analogie restaurant ═══════════
    section_slide(prs, 2, "Comment ca fonctionne ?", "L'analogie du restaurant")

    slide = content_slide(prs, "Imaginez un restaurant...", "Pour comprendre comment l'application est construite")
    add_card(slide, Inches(0.3), Inches(1.7), Inches(3.0), Inches(4.8), "La salle", [
        "C'est ce que les clients", "voient : le menu, les tables,", "la deco, les serveurs.", "",
        "= Le SITE WEB", "(ce que vous voyez quand", "vous ouvrez l'application", "dans votre navigateur)", "",
        "On appelle ca le FRONTEND.",
    ], accent=BLUE_MED)
    add_card(slide, Inches(3.6), Inches(1.7), Inches(3.0), Inches(4.8), "La cuisine", [
        "C'est la ou les plats", "sont prepares. Les clients", "ne la voient jamais.", "",
        "= Le SERVEUR", "(le programme qui traite", "vos demandes : « affiche", "le planning de Lundi »)", "",
        "On appelle ca le BACKEND.",
    ], accent=GREEN)
    add_card(slide, Inches(6.9), Inches(1.7), Inches(3.0), Inches(4.8), "Le frigo / le stock", [
        "C'est la ou sont gardes", "tous les ingredients.", "",
        "= La BASE DE DONNEES", "(la ou sont stockees toutes", "les infos : employes,", "chantiers, planning,", "plans, questions...)", "",
        "Un simple fichier sur le serveur.",
    ], accent=ORANGE)
    add_card(slide, Inches(10.2), Inches(1.7), Inches(3.0), Inches(4.8), "Le fournisseur", [
        "C'est celui qui livre les", "ingredients au restaurant.", "",
        "= MICROSOFT FABRIC", "(le systeme HK qui envoie", "les donnees des filiales,", "employes, chantiers, etc.", "automatiquement chaque", "heure)",
    ], accent=RED)

    slide = content_slide(prs, "Le parcours d'une action", "Quand vous cliquez sur « Voir le planning de Lundi »")
    mono_block(slide, [
        "", "    VOUS (le navigateur)          LE SERVEUR (la cuisine)          LA BASE DE DONNEES",
        "    --------------------          ------------------------          -----------------",
        "", "    1. Vous cliquez sur",
        "       'Planning Lundi'",
        "              |",
        "              v",
        "    2. Le site web envoie   ------>   3. Le serveur recoit",
        "       une demande                       la demande",
        "       (comme commander                        |",
        "        un plat)                                v",
        "                                     4. Il cherche les     ------>  5. La base renvoie",
        "                                        infos necessaires             les donnees",
        "                                              |",
        "                                              v",
        "    7. Le site affiche     <------   6. Le serveur prepare",
        "       le planning                      la reponse et",
        "       de Lundi !                       l'envoie au site",
        "", "    Tout ca se passe en moins d'une seconde !",
    ], top=Inches(1.4))

    # ═══════════ SECTION 3 : Les outils ═══════════
    section_slide(prs, 3, "Les outils utilises", "A quoi sert chaque brique — explique simplement")

    slide = content_slide(prs, "Les outils — Vue d'ensemble", "Chaque outil a un role precis")
    add_bullets(slide, [
        "GitHub : le coffre-fort qui garde tout l'historique du code (comme Google Drive pour du code)",
        "Docker : l'emballage qui garantit que l'application fonctionne partout de la meme facon",
        "Nginx : le portier qui accueille les visiteurs et les dirige au bon endroit",
        "React : l'outil qui construit l'interface visible (boutons, tableaux, pages...)",
        "Fastify : le moteur invisible qui traite les demandes et va chercher les donnees",
        "SQLite : le classeur ou toutes les donnees sont rangees",
        "Microsoft Fabric : le fournisseur de donnees HK",
        "Azure AD : le vigile qui verifie votre identite quand vous vous connectez",
        "Windsurf : l'outil du developpeur pour ecrire le code (avec une IA qui aide)",
    ], top=Inches(1.7), font_size=13, spacing=Pt(5))

    # GitHub
    slide = content_slide(prs, "GitHub — Le coffre-fort du code")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.5), "C'est quoi ?", [
        "Imaginez un Google Drive,", "mais specialise pour le code.", "",
        "Chaque modification est", "enregistree avec un message", "qui dit ce qui a change.", "",
        "On peut revenir en arriere", "a tout moment, comme un", "historique de versions illimite.",
    ], accent=BLUE_MED)
    add_card(slide, Inches(6.5), Inches(1.7), Inches(6.3), Inches(4.5), "A quoi ca sert ici ?", [
        "1. SAUVEGARDER le code en ligne", "   (meme si le PC tombe en panne,", "   rien n'est perdu)", "",
        "2. PARTAGER le code", "   (un nouveau collegue peut recuperer", "   tout le projet en 2 minutes)", "",
        "3. METTRE A JOUR l'application", "   (le serveur va chercher la derniere", "   version du code sur GitHub)",
    ], accent=GREEN)

    # Docker
    slide = content_slide(prs, "Docker — L'emballage universel")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.5), "C'est quoi ?", [
        "Imaginez que vous envoyez", "un colis. Si vous envoyez les", "pieces en vrac, le destinataire", "ne saura pas les assembler.", "",
        "Docker = vous envoyez le", "meuble DEJA MONTE dans", "une boite parfaite.", "",
        "L'application + tout ce qu'il", "lui faut = dans un conteneur.",
    ], accent=GREEN)
    add_card(slide, Inches(6.5), Inches(1.7), Inches(6.3), Inches(4.5), "Pourquoi c'est important ?", [
        "Sans Docker :", "  'Ca marche sur mon PC mais pas", "   sur le serveur !' (classique)", "",
        "Avec Docker :", "  L'application fonctionne exactement", "  pareil partout.", "",
        "Dans notre cas :", "  2 conteneurs : un pour le site web,", "  un pour le moteur + base de donnees.",
    ], accent=BLUE_MED)

    # Windsurf
    slide = content_slide(prs, "Windsurf — L'outil du developpeur")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.5), "C'est quoi ?", [
        "C'est l'outil dans lequel le", "developpeur ecrit le code.", "",
        "Comme Word pour ecrire", "une lettre, Windsurf c'est", "pour ecrire du code.", "",
        "Sa particularite : il a une", "IA integree (Cascade) qui", "aide a ecrire le code !",
    ], accent=ORANGE)
    add_card(slide, Inches(6.5), Inches(1.7), Inches(6.3), Inches(4.5), "Cascade (l'IA)", [
        "Imaginez un assistant a qui", "vous dites en francais :", "",
        "  'Ajoute un bouton qui", "   exporte le planning en PDF'", "",
        "Et il ecrit le code pour vous !", "",
        "L'essentiel de cette application", "a ete developpe comme ca :", "on decrit ce qu'on veut,", "l'IA le programme.",
    ], accent=BLUE_MED)

    # ═══════════ SECTION 4 : Ou vit l'application ═══════════
    section_slide(prs, 4, "Ou vit l'application ?", "Le serveur et comment y acceder")

    slide = content_slide(prs, "L'application tourne sur un serveur dedie")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.8), "Le serveur (la VM)", [
        "C'est un ordinateur qui tourne", "24h/24, 7j/7, sur le reseau CAA.", "",
        "Il n'a pas d'ecran : on le", "controle a distance.", "",
        "Adresse : 192.168.13.51", "Systeme : Debian 13 (Linux)", "",
        "Il heberge les 2 conteneurs", "Docker (le site web + le moteur).",
    ], accent=BLUE_MED)
    add_card(slide, Inches(6.5), Inches(1.7), Inches(6.3), Inches(4.8), "Comment y acceder ?", [
        "EN TANT QU'UTILISATEUR :", "  Ouvrir un navigateur web", "  Taper : https://192.168.13.51", "  Se connecter avec son compte Microsoft", "  C'est tout !", "",
        "EN TANT QU'ADMINISTRATEUR :", "  Ouvrir mRemoteNG (logiciel Windows)", "  Se connecter en SSH au serveur", "  Taper des commandes pour maintenir", "  l'application (voir logs, redemarrer...)",
    ], accent=GREEN)

    slide = content_slide(prs, "Schema simplifie")
    mono_block(slide, [
        "", "    VOTRE PC                                     LE SERVEUR (192.168.13.51)",
        "    --------                                     ---------",
        "", "    Navigateur web                               Systeme : Debian 13 (Linux)",
        "    (Chrome, Edge...)                                    |",
        "          |                                        ------+------",
        "          |   https://192.168.13.51                |           |",
        "          | -------------------------------->  Le PORTIER   Le MOTEUR",
        "          |                                   (Nginx)     (Fastify)",
        "          |                                       |           |",
        "          |   Vous recevez la page                |      La BASE DE",
        "          | <--------------------------------     |      DONNEES",
        "          |                                       |      (SQLite)",
        "                                            Le site web  Toutes vos",
        "                                            (les pages,  donnees",
        "                                             boutons...)",
    ], top=Inches(1.4))

    # ═══════════ SECTION 5 : Donnees ═══════════
    section_slide(prs, 5, "D'ou viennent les donnees ?", "La synchronisation avec HK / Microsoft Fabric")

    slide = content_slide(prs, "Les donnees arrivent automatiquement de HK")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.0), "Le principe", [
        "Toutes les heures, l'application va chercher les donnees a jour dans le systeme HK (via Microsoft Fabric).",
    ], accent=BLUE_MED)
    add_bullets(slide, [
        "Filiales : CAA Guadeloupe, Martinique, etc. — importees automatiquement",
        "Ateliers et equipes : structure organisationnelle importee de HK",
        "Employes : noms, matricules, qualifications — importes et mis a jour",
        "Chantiers : codes, libelles, dates, statuts — importes de HK",
        "OF et articles : ordres de fabrication et leurs articles — importes",
        "Clients et affaires : informations commerciales — importees",
    ], top=Inches(3.0), font_size=14, spacing=Pt(6))

    slide = content_slide(prs, "Ce qui est local (pas dans HK)")
    add_bullets(slide, [
        "Le planning (qui travaille ou, quand) : cree directement dans l'application",
        "Les absences (CP, RTT, maladie) : saisies dans l'application",
        "Les questions et reponses : creees dans l'application",
        "Les plans de fabrication : crees dans l'application",
        "Les pieces jointes (photos, PDF) : uploadees dans l'application",
        "",
        "Ces donnees ne sont PAS ecrasees par la synchronisation.",
        "La sync importe uniquement les donnees de reference (employes, chantiers, OF...).",
    ], top=Inches(1.8), font_size=14, spacing=Pt(6))

    # ═══════════ SECTION 6 : Mise a jour ═══════════
    section_slide(prs, 6, "Comment on met a jour ?", "Le processus de mise a jour de l'application")

    slide = content_slide(prs, "La mise a jour est automatique !")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.2), "Le principe (3 etapes simples)", [
        "1. Le developpeur envoie ses modifications sur GitHub    2. Le serveur detecte les changements    3. L'application se met a jour toute seule",
    ], accent=GREEN)
    add_bullets(slide, [
        "Le serveur verifie GitHub toutes les 5 minutes",
        "S'il detecte du nouveau code, il telecharge la mise a jour",
        "Il reconstruit l'application automatiquement (via Docker)",
        "L'application redémarre avec la nouvelle version",
        "Aucune action necessaire cote utilisateurs !",
        "",
        "Les utilisateurs voient la mise a jour en rafraichissant leur page (F5).",
        "Les donnees (planning, questions, plans...) ne sont JAMAIS perdues lors d'une mise a jour.",
    ], top=Inches(3.1), font_size=14, spacing=Pt(5))

    slide = content_slide(prs, "Schema de la mise a jour")
    mono_block(slide, [
        "",
        "    PC du developpeur                GitHub                  Serveur (VM)",
        "    --------------------             ------                  -----------",
        "",
        "    1. Le dev modifie    -------->  2. Le code est",
        "       le code et                      stocke en ligne",
        "       l'envoie                              |",
        "                                             |  (toutes les 5 min)",
        "                                             v",
        "                                    3. Le serveur detecte",
        "                                       le nouveau code",
        "                                             |",
        "                                             v",
        "                                    4. Il reconstruit     -------->  5. L'application",
        "                                       l'application                    est a jour !",
        "                                       (Docker rebuild)",
        "",
        "    Resultat : le dev envoie son code, tout le reste est automatique.",
    ], top=Inches(1.4))

    # ═══════════ SECTION 7 : Securite ═══════════
    section_slide(prs, 7, "La securite", "Qui peut acceder, et comment c'est protege")

    slide = content_slide(prs, "L'application est securisee")
    add_card(slide, Inches(0.3), Inches(1.7), Inches(3.9), Inches(4.5), "Connexion", [
        "On se connecte avec son", "compte Microsoft pro (@caa).", "",
        "C'est le meme identifiant", "que pour Outlook ou Teams.", "",
        "Pas de mot de passe", "supplementaire a retenir.", "",
        "Seuls les employes CAA", "autorises peuvent acceder.",
    ], accent=BLUE_MED)
    add_card(slide, Inches(4.5), Inches(1.7), Inches(3.9), Inches(4.5), "Protection des donnees", [
        "Connexion chiffree (HTTPS) :", "les donnees sont protegees", "quand elles transitent.", "",
        "Sauvegardes automatiques :", "chaque nuit a 2h du matin,", "une copie de la base est faite.", "",
        "30 jours de sauvegardes", "conserves sur le serveur.",
    ], accent=GREEN)
    add_card(slide, Inches(8.7), Inches(1.7), Inches(3.9), Inches(4.5), "Acces au serveur", [
        "Le serveur n'est accessible", "que sur le reseau interne CAA.", "",
        "L'administration se fait par", "SSH (connexion securisee", "par cle, pas par mot de passe).", "",
        "Les secrets (mots de passe,", "cles) sont stockes dans un", "fichier protege (.env).",
    ], accent=RED)

    # ═══════════ SECTION 8 : Pour aller plus loin ═══════════
    section_slide(prs, 8, "Pour aller plus loin", "Si vous devez intervenir sur l'application")

    slide = content_slide(prs, "Pour le developpeur : les commandes essentielles")
    add_card(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.5), "Sur le serveur (via SSH)", [
        "docker compose ps", "  -> Voir si tout tourne", "",
        "docker compose logs -f api", "  -> Voir ce qui se passe en direct", "",
        "docker compose restart", "  -> Redemarrer l'application", "",
        "./scripts/deploy.sh", "  -> Forcer une mise a jour",
    ], accent=BLUE_MED)
    add_card(slide, Inches(6.5), Inches(1.7), Inches(6.3), Inches(4.5), "Sur le PC de dev", [
        "git clone [url du repo]", "  -> Recuperer tout le code", "",
        "npm run dev", "  -> Lancer le site en local", "",
        "cd api && npm run dev", "  -> Lancer le serveur en local", "",
        "git push", "  -> Envoyer les modifications",
    ], accent=GREEN)

    slide = content_slide(prs, "Les numeros importants")
    add_card(slide, Inches(0.3), Inches(1.7), Inches(3.9), Inches(3.0), "Adresses", [
        "Application :", "https://192.168.13.51", "",
        "Serveur VM :", "192.168.13.51 (SSH)", "",
        "PC dev : 192.168.13.50",
    ], accent=BLUE_MED)
    add_card(slide, Inches(4.5), Inches(1.7), Inches(3.9), Inches(3.0), "Code source", [
        "GitHub :", "github.com/CYBERTESTCAA/", "PLANNING_PREVISIONNEL", "",
        "Branche principale : main",
    ], accent=GREEN)
    add_card(slide, Inches(8.7), Inches(1.7), Inches(3.9), Inches(3.0), "Automatisations", [
        "Sync HK : toutes les heures", "Mise a jour : toutes les 5 min", "Sauvegarde : chaque nuit a 2h", "",
        "Tout est automatique !",
    ], accent=ORANGE)

    # ── FIN ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_DARK)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1), "Merci !", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6), "Des questions ?", font_size=24, color=RGBColor(0xBF,0xDB,0xFE), alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, Inches(4.5), Inches(4.8), Inches(4), Inches(0.05), BLUE_MED)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5), "Projet SUIVI DE CHANTIERS  •  CAA Guadeloupe  •  2025-2026", font_size=14, color=RGBColor(0x93,0xC5,0xFD), alignment=PP_ALIGN.CENTER)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    path = os.path.join(out_dir, "..", "docs", "Presentation_Technique_Simple.pptx")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"Presentation sauvegardee : {path}")
    return path

if __name__ == "__main__":
    build()
